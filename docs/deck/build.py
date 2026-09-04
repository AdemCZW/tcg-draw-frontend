#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
VaultDraw 賣家合作簡報產生器。

產出：
  docs/deck/VaultDraw-賣家合作簡報.pptx
  docs/deck/outline.md          （同一份內容的文字版，給網頁版用）

原則：
  * 每一頁只回答一件事 ——「這對開池的人有什麼好處／擋掉什麼麻煩」。
  * 所有數字都來自原始碼常數，沒有任何推估或行銷數字。
  * 版面文字與 outline.md 由同一份資料產生，兩邊不會走鐘。
"""

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).resolve().parent
ASSETS = HERE / "assets"
OUT_PPTX = HERE / "VaultDraw-賣家合作簡報.pptx"
OUT_MD = HERE / "outline.md"

# ── 色票：直接取自 src/styles/tokens.css ────────────────────────────────
BG = RGBColor(0x0D, 0x0C, 0x0F)
SURFACE = RGBColor(0x17, 0x16, 0x1A)
SURFACE2 = RGBColor(0x1E, 0x1C, 0x22)
LINE = RGBColor(0x30, 0x2D, 0x38)
TEXT = RGBColor(0xF4, 0xF1, 0xEE)
MUTED = RGBColor(0xA4, 0x9B, 0x96)
FAINT = RGBColor(0x6F, 0x68, 0x62)
ACCENT = RGBColor(0xF7, 0x3B, 0x20)
ACCENT_WASH = RGBColor(0x2C, 0x15, 0x12)
OK = RGBColor(0x35, 0xC9, 0x8A)
OK_WASH = RGBColor(0x11, 0x25, 0x1C)
WARN = RGBColor(0xF0, 0xA1, 0x3A)
WARN_WASH = RGBColor(0x2A, 0x21, 0x14)
DANGER = RGBColor(0xFF, 0x5A, 0x4D)
DANGER_WASH = RGBColor(0x2C, 0x16, 0x14)
GOLD = RGBColor(0xF0, 0xBB, 0x52)
LILAC = RGBColor(0xA9, 0x8C, 0xFF)

FONT = "PingFang TC"
MONO = "Menlo"

W, H = Inches(13.333), Inches(7.5)
ML = Inches(0.78)          # 左邊界
MR = Inches(0.78)          # 右邊界
CW = W - ML - MR           # 內容寬度


# ── 低階工具 ──────────────────────────────────────────────────────────
def _set_font(run, size, bold=False, color=TEXT, font=FONT):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.name = font
    f.color.rgb = color
    # python-pptx 只寫 latin；東亞字型要自己補，否則 PowerPoint 會用替代字型
    rpr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rpr.find(qn(tag))
        if el is None:
            el = rpr.makeelement(qn(tag), {})
            rpr.append(el)
        el.set("typeface", font)


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tb


def para(tf, first=False):
    return tf.paragraphs[0] if first else tf.add_paragraph()


def write(tf, chunks, size=15, color=TEXT, bold=False, space_after=6,
          line=1.35, align=PP_ALIGN.LEFT, first=False, font=FONT):
    """chunks: str 或 [(文字, {覆寫}), ...]"""
    p = para(tf, first)
    p.alignment = align
    p.space_after = Pt(space_after)
    p.line_spacing = line
    if isinstance(chunks, str):
        chunks = [(chunks, {})]
    for txt, over in chunks:
        r = p.add_run()
        r.text = txt
        _set_font(r, over.get("size", size), over.get("bold", bold),
                  over.get("color", color), over.get("font", font))
    return p


def card(slide, x, y, w, h, fill=SURFACE, line_color=LINE, radius=0.035, line_w=1.0):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = radius
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    return sh


def pill(slide, x, y, w, h, text, fill, ink, size=11):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.5
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    write(tf, text, size=size, bold=True, color=ink, align=PP_ALIGN.CENTER,
          first=True, space_after=0)
    return sh


def picture(slide, name, x, y, w, h, border=True):
    """等比縮放塞進 (x,y,w,h)，置中。回傳實際擺放的矩形。"""
    path = ASSETS / name
    iw, ih = Image.open(path).size
    scale = min(w / iw, h / ih)
    pw, ph = int(iw * scale), int(ih * scale)
    px, py = int(x + (w - pw) / 2), int(y + (h - ph) / 2)
    pic = slide.shapes.add_picture(str(path), px, py, pw, ph)
    if border:
        pic.line.color.rgb = LINE
        pic.line.width = Pt(0.75)
    return px, py, pw, ph


def caption(slide, x, y, w, text, align=PP_ALIGN.LEFT):
    tb = textbox(slide, x, y, w, Inches(0.24))
    write(tb.text_frame, text, size=10, color=FAINT, align=align, first=True,
          space_after=0)
    return tb


def arrow(slide, x, y, w, h, color=LINE):
    sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    return sh


def rule(slide, x, y, w, color=LINE, thickness=Pt(1)):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Emu(int(thickness)))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    return sh


# ── 頁面骨架 ──────────────────────────────────────────────────────────
class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width, self.prs.slide_height = W, H
        self.blank = self.prs.slide_layouts[6]
        self.outline = []      # (頁碼, kicker, 標題, 副標, [條目])
        self.n = 0

    def new(self, kicker=None, title=None, sub=None, chrome=True):
        self.n += 1
        s = self.prs.slides.add_slide(self.blank)
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG
        bg.line.fill.background()
        bg.shadow.inherit = False
        bg.text_frame.text = ""

        y = Inches(0.52)
        if kicker:
            tb = textbox(s, ML, y, CW, Inches(0.26))
            write(tb.text_frame, kicker, size=12, bold=True, color=ACCENT,
                  first=True, space_after=0)
            y += Inches(0.32)
        if title:
            tb = textbox(s, ML, y, CW, Inches(0.72))
            write(tb.text_frame, title, size=32, bold=True, color=TEXT,
                  line=1.12, first=True, space_after=0)
            y += Inches(0.70)
        if sub:
            tb = textbox(s, ML, y, Inches(11.0), Inches(0.5))
            write(tb.text_frame, sub, size=15, color=MUTED, line=1.35,
                  first=True, space_after=0)
            y += Inches(0.46)

        if chrome:
            rule(s, ML, Inches(6.92), CW, color=RGBColor(0x22, 0x20, 0x29))
            tb = textbox(s, ML, Inches(7.03), Inches(6.0), Inches(0.24))
            write(tb.text_frame, "VaultDraw · 賣家合作說明", size=9,
                  color=FAINT, first=True, space_after=0)
            tb = textbox(s, W - MR - Inches(1.0), Inches(7.03), Inches(1.0), Inches(0.24))
            write(tb.text_frame, f"{self.n:02d}", size=9, color=FAINT,
                  align=PP_ALIGN.RIGHT, first=True, space_after=0)
        return s, y

    def record(self, kicker, title, sub, bullets):
        self.outline.append((self.n, kicker, title, sub, bullets))


d = Deck()
DEMO = "示意畫面：站上目前為示範資料"

# ══════════════════════════════════════════════════════════════════════
# 01 封面
# ══════════════════════════════════════════════════════════════════════
s, _ = d.new(chrome=False)
# 右側放一張裱框的實際畫面。刻意不做全出血漸層 —— 多段半透明矩形疊出來的
# 假漸層在部分渲染器會出現可見的帶狀接縫，裱框反而更乾淨、也更像「這是產品」。
card(s, Inches(7.06), Inches(1.62), Inches(5.62), Inches(4.28),
     fill=SURFACE, line_color=LINE, radius=0.02)
picture(s, "lobby.png", Inches(7.2), Inches(1.76), Inches(5.34), Inches(4.0))
caption(s, Inches(7.06), Inches(5.98), Inches(5.62),
        "大廳 · 示意畫面：站上目前為示範資料")

tb = textbox(s, ML, Inches(1.42), Inches(7.2), Inches(0.6))
p = write(tb.text_frame, [("Vault", {"color": TEXT}), ("Draw", {"color": ACCENT})],
          size=40, bold=True, first=True, space_after=0)

tb = textbox(s, ML, Inches(2.46), Inches(7.0), Inches(1.9))
write(tb.text_frame, "在這裡開池，", size=44, bold=True, line=1.18, first=True,
      space_after=2)
write(tb.text_frame, [("抽成 0%", {"color": ACCENT}), ("，", {}),
                      ("逐筆結算", {"color": ACCENT})], size=44, bold=True,
      line=1.18, space_after=2)

tb = textbox(s, ML, Inches(4.42), Inches(5.9), Inches(1.2))
write(tb.text_frame,
      "給台灣卡舖、卡商與個人大盤的合作說明。卡從頭到尾在你手上，"
      "平台只負責制度：公平性、金流保管、糾紛的客觀答案。",
      size=15, color=MUTED, line=1.55, first=True, space_after=0)

for i, t in enumerate(["平台抽成 0%", "買回價自己定", "一卡不會被賣兩次"]):
    pill(s, ML + Inches(i * 2.06), Inches(5.72), Inches(1.9), Inches(0.42),
         t, SURFACE2, TEXT, size=11.5)

caption(s, ML, Inches(6.55), Inches(6.0), "本簡報所有數字皆取自系統設定值 · 平台尚未正式營運")
d.record("封面", "在這裡開池，抽成 0%，逐筆結算",
         "給台灣卡舖、卡商與個人大盤的合作說明。",
         ["卡從頭到尾在你手上，平台只負責制度：公平性、金流保管、糾紛的客觀答案。",
          "三個重點：平台抽成 0% / 買回價自己定 / 一卡不會被賣兩次。",
          "本簡報所有數字皆取自系統設定值；平台尚未正式營運。"])

# ══════════════════════════════════════════════════════════════════════
# 02 一句話講完
# ══════════════════════════════════════════════════════════════════════
s, y = d.new("一句話", "你出卡，我們出制度。",
             "VaultDraw 是鑑定卡的線上抽選（オリパ）平台。你把手上的卡編成一個池、"
             "自己定價、自己宣告買回價；平台處理籤序、金流與交付的規則。")

items = [
    ("卡在你手上", "平台不代管實體卡。買家抽中後，你直接寄給他 —— "
                   "平台只提供收件資訊與一個雙方確認的機制。", GOLD),
    ("錢在制度裡", "票金先保管，逐筆釋放。你不用等整池抽完才拿得到錢，"
                   "買家也不用擔心付了錢拿不到卡。", OK),
    ("公平性可驗算", "籤序在開賣前就封存並公布雜湊，完抽後公開種子。"
                     "「是不是有內定」這件事有客觀答案，不需要你自證清白。", LILAC),
]
cw = (CW - Inches(0.44)) / 3
for i, (t, body, c) in enumerate(items):
    x = ML + int(i * (cw + Inches(0.22)))
    card(s, x, Inches(2.62), int(cw), Inches(2.6))
    rule(s, x + Inches(0.34), Inches(2.62), Inches(0.9), color=c, thickness=Pt(3))
    tb = textbox(s, x + Inches(0.34), Inches(3.05), int(cw) - Inches(0.68), Inches(0.4))
    write(tb.text_frame, t, size=19, bold=True, color=TEXT, first=True, space_after=0)
    tb = textbox(s, x + Inches(0.34), Inches(3.62), int(cw) - Inches(0.68), Inches(1.4))
    write(tb.text_frame, body, size=13, color=MUTED, line=1.5, first=True, space_after=0)

tb = textbox(s, ML, Inches(5.52), CW, Inches(0.9))
write(tb.text_frame,
      [("換句話說：", {"color": MUTED}),
       ("你原本花在自證清白、對帳、處理糾紛上的時間，這裡由系統的規則承擔。",
        {"color": TEXT, "bold": True})],
      size=16, line=1.4, first=True, space_after=0)
d.record("一句話", "你出卡，我們出制度。",
         "VaultDraw 是鑑定卡的線上抽選（オリパ）平台。你把手上的卡編成一個池、自己定價、自己宣告買回價；平台處理籤序、金流與交付的規則。",
         ["卡在你手上 —— 平台不代管實體卡。買家抽中後你直接寄給他，平台只提供收件資訊與一個雙方確認的機制。",
          "錢在制度裡 —— 票金先保管、逐筆釋放。你不用等整池抽完才拿得到錢。",
          "公平性可驗算 —— 籤序開賣前封存並公布雜湊，完抽後公開種子。「是不是有內定」有客觀答案。",
          "換句話說：你原本花在自證清白、對帳、處理糾紛上的時間，這裡由系統的規則承擔。"])

# ══════════════════════════════════════════════════════════════════════
# 03 痛點
# ══════════════════════════════════════════════════════════════════════
s, y = d.new("你現在在處理的四件事", "這四件麻煩，沒有一件跟「賣卡」有關。",
             "不管你現在在別的オリパ平台開池，還是在蝦皮／FB 社團賣卡，時間大多花在這裡。")

pains = [
    ("「你們是不是有內定？」",
     "抽到爛卡的人在留言區問，你拿不出證據，只能一個一個解釋。解釋一百次，"
     "還是有人不信。"),
    ("同一張卡被登記兩次",
     "上架時貼錯編號、或同一張卡同時掛在兩個地方 —— 通常是兩個買家吵起來才發現，"
     "傷的是你的信譽。"),
    ("錢卡在別人手上",
     "整池賣完才結算，或是等平台的月結。現金流被別人的節奏綁住，"
     "下一批貨進不了。"),
    ("客服與對帳",
     "誰付了錢、誰還沒寄、哪一筆該退 —— 靠對話紀錄與試算表追。"
     "規模一大就開始漏。"),
]
cw2 = (CW - Inches(0.3)) / 2
for i, (t, body) in enumerate(pains):
    col, row = i % 2, i // 2
    x = ML + int(col * (cw2 + Inches(0.3)))
    yy = Inches(2.48) + int(row * Inches(1.92))
    card(s, x, yy, int(cw2), Inches(1.68), fill=SURFACE, line_color=LINE)
    pill(s, x + Inches(0.3), yy + Inches(0.28), Inches(0.32), Inches(0.32),
         "✕", DANGER_WASH, DANGER, size=12)
    tb = textbox(s, x + Inches(0.76), yy + Inches(0.26), int(cw2) - Inches(1.1), Inches(0.36))
    write(tb.text_frame, t, size=16, bold=True, color=TEXT, first=True, space_after=0)
    tb = textbox(s, x + Inches(0.76), yy + Inches(0.72), int(cw2) - Inches(1.1), Inches(0.86))
    write(tb.text_frame, body, size=12.5, color=MUTED, line=1.45, first=True, space_after=0)

tb = textbox(s, ML, Inches(6.34), CW, Inches(0.4))
write(tb.text_frame,
      [("接下來四頁，", {"color": MUTED}),
       ("每一頁對應上面一件事，講的是系統怎麼替你擋掉它。", {"color": TEXT, "bold": True})],
      size=14, first=True, space_after=0)
d.record("你現在在處理的四件事", "這四件麻煩，沒有一件跟「賣卡」有關。",
         "不管你現在在別的オリパ平台開池，還是在蝦皮／FB 社團賣卡，時間大多花在這裡。",
         ["「你們是不是有內定？」—— 抽到爛卡的人在留言區問，你拿不出證據，只能一個一個解釋。",
          "同一張卡被登記兩次 —— 通常是兩個買家吵起來才發現，傷的是你的信譽。",
          "錢卡在別人手上 —— 整池賣完才結算或等月結，現金流被別人的節奏綁住。",
          "客服與對帳 —— 誰付了錢、誰還沒寄、哪一筆該退，靠對話紀錄與試算表追。",
          "接下來四頁，每一頁對應上面一件事，講的是系統怎麼替你擋掉它。"])

# ══════════════════════════════════════════════════════════════════════
# 04 抽成 0%
# ══════════════════════════════════════════════════════════════════════
s, y = d.new("抽成", "平台抽成 0%。",
             "不是首年優惠、不是階梯費率、沒有上架費。系統裡那個費率常數的值就是 0。")

card(s, ML, Inches(2.42), Inches(5.6), Inches(3.0), fill=ACCENT_WASH, line_color=ACCENT)
tb = textbox(s, ML + Inches(0.5), Inches(2.86), Inches(4.6), Inches(1.5))
write(tb.text_frame, "0%", size=96, bold=True, color=ACCENT, line=1.0, first=True,
      space_after=0)
tb = textbox(s, ML + Inches(0.5), Inches(4.32), Inches(4.6), Inches(0.9))
write(tb.text_frame, "平台從你的票收抽走的比例", size=15, color=TEXT, first=True,
      space_after=4)
write(tb.text_frame, "PLATFORM_FEE_RATE = 0", size=12, color=MUTED, font=MONO,
      space_after=0)

rx = ML + Inches(6.0)
rw = CW - Inches(6.0)
notes = [
    ("票收 100% 是你的",
     "買家付的每一點票金，扣掉你自己宣告的買回價之後，剩下的全部進你的口袋。"),
    ("沒有上架費、沒有開池費",
     "開一個池不用先付錢。第一個池賣不好，你的損失是時間，不是費用。"),
    ("這一頁刻意不談我們怎麼賺錢",
     "平台還在測試階段，收入模式尚未定案。定案之前的規則就是這一頁寫的，"
     "改動會提前公告，不會回頭對已開的池追收。"),
]
yy = Inches(2.42)
for t, body in notes:
    tb = textbox(s, rx, yy, rw, Inches(0.32))
    write(tb.text_frame, t, size=16, bold=True, color=TEXT, first=True, space_after=0)
    tb = textbox(s, rx, yy + Inches(0.38), rw, Inches(0.72))
    write(tb.text_frame, body, size=13, color=MUTED, line=1.5, first=True, space_after=0)
    yy += Inches(1.06)

d.record("抽成", "平台抽成 0%。",
         "不是首年優惠、不是階梯費率、沒有上架費。系統裡那個費率常數的值就是 0（PLATFORM_FEE_RATE = 0）。",
         ["票收 100% 是你的 —— 買家付的每一點票金，扣掉你自己宣告的買回價之後，剩下的全部進你的口袋。",
          "沒有上架費、沒有開池費 —— 開一個池不用先付錢。第一個池賣不好，損失是時間不是費用。",
          "這一頁刻意不談我們怎麼賺錢 —— 平台還在測試階段，收入模式尚未定案。改動會提前公告，不會回頭對已開的池追收。"])

# ══════════════════════════════════════════════════════════════════════
# 05 錢什麼時候進口袋（時間軸）
# ══════════════════════════════════════════════════════════════════════
s, y = d.new("現金流", "逐筆結算，不用等整池抽完。",
             "一張卡被抽走，那一筆票金就記在你名下。它走完交付流程就變成可動用 —— "
             "不必等同池的其他籤賣掉。")

TY = Inches(2.72)          # 時間軸基準線
steps = [
    ("抽中", "票金記在你名下\n（保留額）", ACCENT),
    ("你寄出", "72 小時內\n寄給買家", GOLD),
    ("買家確認／7 天期滿", "鑑賞期走完\n（以先到者為準）", OK),
    ("可動用", "那一筆入袋", OK),
]
sw = Inches(2.42)
gap = Inches(0.62)
x0 = ML + Inches(0.1)
rule(s, x0 + Inches(0.2), TY + Inches(0.62), Inches(10.9), color=RGBColor(0x2A, 0x27, 0x32),
     thickness=Pt(2))
for i, (t, body, c) in enumerate(steps):
    x = x0 + int(i * (sw + gap))
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.06), TY + Inches(0.44),
                             Inches(0.36), Inches(0.36))
    dot.fill.solid(); dot.fill.fore_color.rgb = c
    dot.line.color.rgb = BG; dot.line.width = Pt(3)
    dot.shadow.inherit = False; dot.text_frame.text = ""
    tb = textbox(s, x, TY, sw, Inches(0.44))
    write(tb.text_frame, t, size=15, bold=True, color=TEXT, line=1.15, first=True,
          space_after=0)
    tb = textbox(s, x, TY + Inches(0.98), sw, Inches(0.7))
    write(tb.text_frame, body, size=12, color=MUTED, line=1.35, first=True, space_after=0)
    if i < len(steps) - 1:
        arrow(s, x + sw + Inches(0.12), TY + Inches(0.54), Inches(0.38), Inches(0.16),
              color=RGBColor(0x4A, 0x45, 0x55))

# 例外分支
bx, by = ML, Inches(4.5)
card(s, bx, by, Inches(6.2), Inches(1.72), fill=SURFACE, line_color=LINE)
pill(s, bx + Inches(0.32), by + Inches(0.3), Inches(1.42), Inches(0.32),
     "買家不申請出貨", SURFACE2, GOLD, size=11)
tb = textbox(s, bx + Inches(0.32), by + Inches(0.78), Inches(5.56), Inches(0.8))
write(tb.text_frame,
      [("14 天", {"bold": True, "color": GOLD}),
       ("後視同接受寄存，那一筆自動釋放。你的出貨義務不會消失 —— "
        "買家之後仍可申請 —— 但錢不再被別人的沉默綁住。", {})],
      size=13, color=MUTED, line=1.5, first=True, space_after=0)

card(s, bx + Inches(6.5), by, Inches(5.27), Inches(1.72), fill=WARN_WASH, line_color=WARN)
pill(s, bx + Inches(6.82), by + Inches(0.3), Inches(1.0), Inches(0.32),
     "逾期未寄", WARN_WASH, WARN, size=11)
tb = textbox(s, bx + Inches(6.82), by + Inches(0.78), Inches(4.63), Inches(0.8))
write(tb.text_frame,
      [("超過 72 小時", {"bold": True, "color": WARN}),
       ("，系統把票金退還買家，並記你一次違約。這條對你是限制，"
        "對整個站的信任是必要的。", {})],
      size=13, color=MUTED, line=1.5, first=True, space_after=0)

d.record("現金流", "逐筆結算，不用等整池抽完。",
         "一張卡被抽走，那一筆票金就記在你名下。它走完交付流程就變成可動用 —— 不必等同池的其他籤賣掉。",
         ["時間軸：抽中（票金記在你名下・保留額）→ 你寄出（72 小時內）→ 買家確認收貨或 7 天鑑賞期滿 → 那一筆變成可動用。",
          "買家一直不申請出貨：14 天後視同接受寄存，那一筆自動釋放。你的出貨義務不會消失，但錢不再被別人的沉默綁住。",
          "逾期未寄：超過 72 小時，系統把票金退還買家，並記你一次違約。這條對你是限制，對整個站的信任是必要的。"])

# ══════════════════════════════════════════════════════════════════════
# 06 買回價自己定
# ══════════════════════════════════════════════════════════════════════
s, y = d.new("定價權", "買回價是你自己宣告的金額，不是平台訂的比率。",
             "一個賞別填一個絕對金額（例如 A 賞 3,000、D 賞 120），範圍 10 – 10,000,000 點。"
             "沒有基準、沒有係數、平台不會事後調。")

picture(s, "x-tierbuy.png", ML, Inches(2.62), Inches(7.1), Inches(1.9))
caption(s, ML, Inches(4.58), Inches(7.1), f"開池表單的買回價欄位 · {DEMO}")
picture(s, "x-ret.png", ML, Inches(4.98), Inches(7.1), Inches(1.3))
caption(s, ML, Inches(6.36), Inches(7.1), f"買家在池頁看到的同一組數字 · {DEMO}")

rx = ML + Inches(7.5)
rw = CW - Inches(7.5)
card(s, rx, Inches(2.62), rw, Inches(3.74), fill=SURFACE, line_color=LINE)
tb = textbox(s, rx + Inches(0.34), Inches(2.94), rw - Inches(0.68), Inches(0.34))
write(tb.text_frame, "為什麼這對你有利", size=16, bold=True, color=TEXT, first=True,
      space_after=0)
pts = [
    ("你說多少就是多少", "沒有「平台參考價的 X 成」這種會被事後調整的公式。"),
    ("開賣後改不了", "金額被寫進公開的公平性承諾雜湊。這一池的成本上限當場鎖死。"),
    ("下限不是平均", "保底回饋率算的是最壞情況（整池被全數買回），"
                     "不是你預期要付的錢。"),
]
yy = Inches(3.44)
for t, b in pts:
    tb = textbox(s, rx + Inches(0.34), yy, rw - Inches(0.68), Inches(0.3))
    write(tb.text_frame, [("· ", {"color": ACCENT}), (t, {})], size=14, bold=True,
          color=TEXT, first=True, space_after=0)
    tb = textbox(s, rx + Inches(0.52), yy + Inches(0.32), rw - Inches(0.86), Inches(0.66))
    write(tb.text_frame, b, size=12, color=MUTED, line=1.45, first=True, space_after=0)
    yy += Inches(1.0)

d.record("定價權", "買回價是你自己宣告的金額，不是平台訂的比率。",
         "一個賞別填一個絕對金額（例如 A 賞 3,000、D 賞 120），範圍 10 – 10,000,000 點。沒有基準、沒有係數、平台不會事後調。",
         ["你說多少就是多少 —— 沒有「平台參考價的 X 成」這種會被事後調整的公式。",
          "開賣後改不了 —— 那個金額被寫進公開的公平性承諾雜湊。對買家是保障，對你是「這一池的成本上限當場鎖死」。",
          "下限不是平均 —— 保底回饋率＝Σ(買回價×數量) ÷ 票收，算的是最壞情況，不是你預期要付的錢。",
          "畫面：開池表單的買回價欄位（依賞別）、以及買家在池頁看到的同一組數字。"])

# ══════════════════════════════════════════════════════════════════════
# 07 一卡不會被賣兩次
# ══════════════════════════════════════════════════════════════════════
s, y = d.new("信譽保護", "同一個鑑定編號，不可能出現在第二個池。",
             "鑑定編號在資料庫層是唯一鍵 unique(grader, cert_no)。"
             "不是「上架時提醒一下」，是寫進資料庫、繞不過去的約束。")

# 三段防線
lanes = [
    ("你按下開池", "系統把每一張卡的\n鑑定編號一起送出", SURFACE2, MUTED),
    ("資料庫唯一鍵", "unique(grader, cert_no)\n第二次寫入直接失敗", ACCENT_WASH, ACCENT),
    ("當場被擋下", "不是等兩個買家\n吵起來才發現", OK_WASH, OK),
]
lw = Inches(3.1)
for i, (t, b, fill, ink) in enumerate(lanes):
    x = ML + int(i * (lw + Inches(0.52)))
    card(s, x, Inches(2.66), lw, Inches(1.62), fill=fill, line_color=ink)
    tb = textbox(s, x + Inches(0.28), Inches(2.94), lw - Inches(0.56), Inches(0.32))
    write(tb.text_frame, t, size=15, bold=True, color=ink if i else TEXT, first=True,
          space_after=0)
    tb = textbox(s, x + Inches(0.28), Inches(3.36), lw - Inches(0.56), Inches(0.7))
    write(tb.text_frame, b, size=12, color=MUTED, line=1.4, first=True, space_after=0)
    if i < 2:
        arrow(s, x + lw + Inches(0.06), Inches(3.4), Inches(0.4), Inches(0.16),
              color=RGBColor(0x4A, 0x45, 0x55))

picture(s, "x-certdup.png", ML, Inches(4.62), Inches(11.0), Inches(1.06))
caption(s, ML, Inches(5.74), Inches(11.0),
        f"從既有的池複製範本時，鑑定卡那一列會被擋下並要求重挑 · {DEMO}")

tb = textbox(s, ML, Inches(6.16), CW, Inches(0.6))
write(tb.text_frame,
      [("這條保護的是你的信譽。", {"bold": True, "color": TEXT}),
       ("買家不需要相信你有好好管庫存 —— 系統本身就不允許那個錯誤發生。"
        "沒有編號的生卡（RAW）不受此約束，但它們在池頁會被標成「未鑑定」，"
        "買家看得到。", {"color": MUTED})],
      size=13, line=1.5, first=True, space_after=0)

d.record("信譽保護", "同一個鑑定編號，不可能出現在第二個池。",
         "鑑定編號在資料庫層是唯一鍵 unique(grader, cert_no)。不是「上架時提醒一下」，是寫進資料庫、繞不過去的約束。",
         ["你按下開池 → 系統把每一張卡的鑑定編號一起送出。",
          "資料庫唯一鍵 unique(grader, cert_no) → 第二次寫入直接失敗。",
          "當場被擋下 —— 不是等兩個買家吵起來才發現。",
          "這條保護的是你的信譽：買家不需要相信你有好好管庫存，系統本身就不允許那個錯誤發生。",
          "沒有編號的生卡（RAW）不受此約束，但在池頁會被標成「未鑑定」，買家看得到。"])

# ══════════════════════════════════════════════════════════════════════
# 08 公平性可驗算
# ══════════════════════════════════════════════════════════════════════
s, y = d.new("客訴的客觀答案", "「是不是有內定」這件事，有答案可以算。",
             "籤序在開賣前就洗好、封存，並公布一組 SHA-256 承諾雜湊。"
             "隨機來源用 drand（公開的分散式亂數）的未來輪次 —— 那個數字在開池當下還不存在。")

phases = [
    ("1", "開池前 · 承諾", "產生 server seed，只公布它的 SHA-256。"
                           "籤序由這組 seed 決定，公布之後改不動而不被發現。", LILAC),
    ("2", "開賣時 · 鎖定外部亂數", "鎖定 drand 一個「還沒發生」的未來輪次。"
                                   "那一輪的值全世界都還不知道，包括平台 —— "
                                   "所以沒有人能試到對自己有利的籤序。", ACCENT),
    ("3", "完抽後 · 揭示", "公開 server seed 與完整籤序。任何人都能自己重算一次，"
                           "先驗雜湊、再對籤序。", OK),
]
pw3 = (CW - Inches(0.44)) / 3
for i, (num, t, b, c) in enumerate(phases):
    x = ML + int(i * (pw3 + Inches(0.22)))
    card(s, x, Inches(2.78), int(pw3), Inches(2.02))
    pill(s, x + Inches(0.3), Inches(3.04), Inches(0.34), Inches(0.34), num,
         SURFACE2, c, size=13)
    tb = textbox(s, x + Inches(0.76), Inches(3.06), int(pw3) - Inches(1.06), Inches(0.32))
    write(tb.text_frame, t, size=14.5, bold=True, color=c, first=True, space_after=0)
    tb = textbox(s, x + Inches(0.3), Inches(3.56), int(pw3) - Inches(0.6), Inches(1.1))
    write(tb.text_frame, b, size=12, color=MUTED, line=1.45, first=True, space_after=0)
    if i < 2:
        arrow(s, x + int(pw3) + Inches(0.02), Inches(3.7), Inches(0.18), Inches(0.14),
              color=RGBColor(0x4A, 0x45, 0x55))

picture(s, "x-fair.png", ML, Inches(5.06), Inches(5.4), Inches(1.24))
caption(s, ML, Inches(6.36), Inches(5.4), f"每一池的驗證分頁 · {DEMO}")

rx = ML + Inches(5.9)
card(s, rx, Inches(5.06), CW - Inches(5.9), Inches(1.24), fill=SURFACE, line_color=LINE)
tb = textbox(s, rx + Inches(0.32), Inches(5.3), CW - Inches(6.54), Inches(0.8))
write(tb.text_frame,
      [("對你的意義：", {"bold": True, "color": TEXT}),
       ("下次有人在留言區問「是不是有內定」，你不用寫一段解釋 —— "
        "把驗證分頁的連結貼給他，讓他自己算。", {"color": MUTED})],
      size=13, line=1.5, first=True, space_after=0)

d.record("客訴的客觀答案", "「是不是有內定」這件事，有答案可以算。",
         "籤序在開賣前就洗好、封存，並公布一組 SHA-256 承諾雜湊。隨機來源用 drand（公開的分散式亂數）的未來輪次 —— 那個數字在開池當下還不存在。",
         ["1 開池前・承諾：產生 server seed，只公布它的 SHA-256。籤序由這組 seed 決定，公布之後改不動而不被發現。",
          "2 開賣時・鎖定外部亂數：鎖定 drand 一個「還沒發生」的未來輪次。那一輪的值全世界都還不知道，包括平台 —— 所以沒有人能試到對自己有利的籤序。",
          "3 完抽後・揭示：公開 server seed 與完整籤序。任何人都能自己重算一次，先驗雜湊、再對籤序。",
          "對你的意義：下次有人問「是不是有內定」，你不用寫一段解釋 —— 把驗證分頁的連結貼給他，讓他自己算。"])

# ══════════════════════════════════════════════════════════════════════
# 09 開池流程
# ══════════════════════════════════════════════════════════════════════
s, y = d.new("開一個池", "一頁表單，右側即時試算，缺什麼直接告訴你。",
             "沒有審核排隊、沒有客服來回。表單填完就能開賣。")

picture(s, "seller-new-filled-top.png", ML, Inches(2.5), Inches(7.5), Inches(3.9))
caption(s, ML, Inches(6.46), Inches(7.5), f"開池表單 · {DEMO}")

rx = ML + Inches(7.86)
rw = CW - Inches(7.86)
flow = [
    ("挑卡", "從你的卡冊挑，或搜卡片目錄。卡號、系列、版本一起帶進獎品。"),
    ("配賞別與數量", "A/B/C/D 賞與最後賞，各自填數量。"),
    ("填票價與買回價", "右側即時算出票收與保底回饋率。"),
    ("按下開池", "系統當場洗籤、封存、公布雜湊，池立刻上架。"),
]
yy = Inches(2.5)
for i, (t, b) in enumerate(flow):
    pill(s, rx, yy + Inches(0.02), Inches(0.3), Inches(0.3), str(i + 1),
         ACCENT_WASH, ACCENT, size=12)
    tb = textbox(s, rx + Inches(0.44), yy, rw - Inches(0.44), Inches(0.3))
    write(tb.text_frame, t, size=14.5, bold=True, color=TEXT, first=True, space_after=0)
    tb = textbox(s, rx + Inches(0.44), yy + Inches(0.32), rw - Inches(0.44), Inches(0.66))
    write(tb.text_frame, b, size=12, color=MUTED, line=1.42, first=True, space_after=0)
    yy += Inches(0.95)

card(s, rx, Inches(6.34), rw, Inches(0.46), fill=SURFACE2, line_color=None)
tb = textbox(s, rx + Inches(0.22), Inches(6.46), rw - Inches(0.44), Inches(0.3))
write(tb.text_frame,
      [("販售天數自訂：", {"bold": True, "color": TEXT}),
       ("預設 14 天，最長 90 天。", {"color": MUTED})],
      size=12, first=True, space_after=0)

d.record("開一個池", "一頁表單，右側即時試算，缺什麼直接告訴你。",
         "沒有審核排隊、沒有客服來回。表單填完就能開賣。",
         ["1 挑卡 —— 從你的卡冊挑，或搜卡片目錄。卡號、系列、版本一起帶進獎品。",
          "2 配賞別與數量 —— A/B/C/D 賞與最後賞，各自填數量。",
          "3 填票價與買回價 —— 右側即時算出票收與保底回饋率。",
          "4 按下開池 —— 系統當場洗籤、封存、公布雜湊，池立刻上架。",
          "販售天數自訂：預設 14 天，最長 90 天。",
          "也可以「以既有的池為範本」一次帶進卡片、賞別、數量、票價與買回價（籤序與種子一定重新產生）。"])

# ══════════════════════════════════════════════════════════════════════
# 10 經濟護欄
# ══════════════════════════════════════════════════════════════════════
s, y = d.new("經濟護欄", "建池時先算一次最壞情況，算不過就不讓你開。",
             "保底回饋率 ＝ Σ(買回價 × 數量) ÷ 票收。這是限制，也是保護 —— "
             "它擋掉的是「當下沒算清楚、上線之後才發現在倒貼」。")

guards = [
    ("x-econ-over.png", "100% 以上　擋下", "整池抽完再全部買回會倒貼。那不是促銷，是印鈔機。", DANGER),
    ("x-econ-warn.png", "90% 以上　提醒", "放行，但告訴你扣掉運費與金流成本可能倒貼。", WARN),
    ("x-econ-under.png", "低於 25%　擋下", "保底形同沒有，對玩家不公平。站的信任撐不住這種池。", DANGER),
    ("x-econ-ok.png", "合理區間　放行", "落在區間內就直接開，不需要任何人審。", OK),
]
gw = (CW - Inches(0.54)) / 4
for i, (img, t, b, c) in enumerate(guards):
    x = ML + int(i * (gw + Inches(0.18)))
    card(s, x, Inches(2.68), int(gw), Inches(3.30), fill=SURFACE, line_color=LINE)
    picture(s, img, x + Inches(0.16), Inches(2.84), int(gw) - Inches(0.32), Inches(1.72))
    tb = textbox(s, x + Inches(0.24), Inches(4.72), int(gw) - Inches(0.48), Inches(0.3))
    write(tb.text_frame, t, size=14, bold=True, color=c, first=True, space_after=0)
    tb = textbox(s, x + Inches(0.24), Inches(5.12), int(gw) - Inches(0.48), Inches(1.0))
    write(tb.text_frame, b, size=11.5, color=MUTED, line=1.42, first=True, space_after=0)

caption(s, ML, Inches(6.44), CW, f"開池表單右側的即時試算，四種狀態 · {DEMO}")

d.record("經濟護欄", "建池時先算一次最壞情況，算不過就不讓你開。",
         "保底回饋率 ＝ Σ(買回價 × 數量) ÷ 票收。這是限制，也是保護 —— 它擋掉的是「當下沒算清楚、上線之後才發現在倒貼」。",
         ["100% 以上　擋下 —— 整池抽完再全部買回會倒貼。那不是促銷，是印鈔機。",
          "90% 以上　提醒 —— 放行，但告訴你扣掉運費與金流成本可能倒貼。",
          "低於 25%　擋下 —— 保底形同沒有，對玩家不公平。站的信任撐不住這種池。",
          "落在合理區間　放行 —— 直接開，不需要任何人審。"])

# ══════════════════════════════════════════════════════════════════════
# 11 出貨與結算
# ══════════════════════════════════════════════════════════════════════
s, y = d.new("出貨與結算", "哪一筆該寄、哪一筆已入袋，同一頁看完。",
             "不用自己開試算表對帳。錢分成「可動用」與「保留中」，"
             "保留中的每一筆都標著它卡在哪一關。")

picture(s, "x-money.png", ML, Inches(2.56), Inches(7.3), Inches(2.14))
caption(s, ML, Inches(4.78), Inches(7.3), f"保留額依關卡分成三堆 · {DEMO}")

picture(s, "x-shiprow.png", ML, Inches(5.06), Inches(6.9), Inches(1.48))
caption(s, ML, Inches(6.62), Inches(6.9), "每一筆都有倒數與一鍵標記 · 同一個買家的多張卡可以一次標記")

rx = ML + Inches(7.66)
rw = CW - Inches(7.66)
card(s, rx, Inches(2.56), rw, Inches(4.2), fill=SURFACE, line_color=LINE)
tb = textbox(s, rx + Inches(0.3), Inches(2.84), rw - Inches(0.6), Inches(0.32))
write(tb.text_frame, "平台不代管實體卡", size=16, bold=True, color=TEXT, first=True,
      space_after=0)
tb = textbox(s, rx + Inches(0.3), Inches(3.26), rw - Inches(0.6), Inches(1.0))
write(tb.text_frame,
      "卡從頭到尾在你手上，直接寄給買家。平台提供收件資訊與雙方確認的機制，"
      "不插手寄送、不收保管費。",
      size=12.5, color=MUTED, line=1.5, first=True, space_after=0)
rule(s, rx + Inches(0.3), Inches(4.30), rw - Inches(0.6))
rows = [
    ("等你寄出", "72 小時內寄出並標記", GOLD),
    ("寄出後鑑賞期", "買家確認或 7 天期滿", OK),
    ("買家還沒申請", "14 天後自動釋放", MUTED),
]
yy = Inches(4.52)
for t, b, c in rows:
    tb = textbox(s, rx + Inches(0.3), yy, Inches(1.5), Inches(0.28))
    write(tb.text_frame, t, size=12.5, bold=True, color=c, first=True, space_after=0)
    tb = textbox(s, rx + Inches(1.86), yy, rw - Inches(2.16), Inches(0.28))
    write(tb.text_frame, b, size=12.5, color=MUTED, align=PP_ALIGN.RIGHT, first=True,
          space_after=0)
    yy += Inches(0.44)
tb = textbox(s, rx + Inches(0.3), Inches(5.94), rw - Inches(0.6), Inches(0.7))
write(tb.text_frame,
      "違約（逾期未出貨）累積 3 次就不能再開新池。沒有保證金，這是唯一擋得住"
      "連續違約的手段。",
      size=11.5, color=FAINT, line=1.45, first=True, space_after=0)

d.record("出貨與結算", "哪一筆該寄、哪一筆已入袋，同一頁看完。",
         "不用自己開試算表對帳。錢分成「可動用」與「保留中」，保留中的每一筆都標著它卡在哪一關。",
         ["保留額依關卡分成三堆：等你寄出（72 小時內寄出並標記）／寄出後鑑賞期（買家確認或 7 天期滿）／買家還沒申請（14 天後自動釋放）。",
          "每一筆都有倒數與一鍵標記已出貨；同一個買家的多張卡可以裝一起寄、一次標記。",
          "平台不代管實體卡 —— 卡從頭到尾在你手上，直接寄給買家。平台提供收件資訊與雙方確認的機制，不插手寄送、不收保管費。",
          "違約（逾期未出貨）累積 3 次，就不能再開新池。沒有保證金，這是唯一擋得住連續違約的手段。"])

# ══════════════════════════════════════════════════════════════════════
# 12 新賣家怎麼開始
# ══════════════════════════════════════════════════════════════════════
s, y = d.new("新賣家", "第一個池有上限。這是沒有保證金的代價，也是它的好處。",
             "我們不收保證金。取而代之的是把第一次違約的最大損失壓住 —— "
             "限制只套用在你的第一個池。")

card(s, ML, Inches(2.62), Inches(5.5), Inches(2.5), fill=WARN_WASH, line_color=WARN)
tb = textbox(s, ML + Inches(0.4), Inches(2.94), Inches(4.7), Inches(0.3))
write(tb.text_frame, "第一個池的上限", size=14, bold=True, color=WARN, first=True,
      space_after=0)
tb = textbox(s, ML + Inches(0.4), Inches(3.36), Inches(2.2), Inches(0.9))
write(tb.text_frame, "100", size=44, bold=True, color=TEXT, line=1.0, first=True,
      space_after=2)
write(tb.text_frame, "籤", size=13, color=MUTED, space_after=0)
tb = textbox(s, ML + Inches(2.7), Inches(3.36), Inches(2.4), Inches(0.9))
write(tb.text_frame, "100,000", size=44, bold=True, color=TEXT, line=1.0, first=True,
      space_after=2)
write(tb.text_frame, "點票收", size=13, color=MUTED, space_after=0)
tb = textbox(s, ML + Inches(0.4), Inches(4.5), Inches(4.7), Inches(0.4))
write(tb.text_frame,
      [("完成第一個池就解除。", {"bold": True, "color": TEXT}),
       ("之後沒有任何規模限制。", {"color": MUTED})],
      size=13, first=True, space_after=0)

rx = ML + Inches(5.9)
rw = CW - Inches(5.9)
tb = textbox(s, rx, Inches(2.62), rw, Inches(0.3))
write(tb.text_frame, "開始要準備的東西", size=16, bold=True, color=TEXT, first=True,
      space_after=0)
prep = [
    ("卡", "手上要有實體卡。鑑定卡填鑑定機構與編號，生卡直接登記。"),
    ("收件資訊", "買家申請出貨時，平台把地址給你，你直接寄。"),
    ("一個帳號", "LINE 或 Email 登入即可，不用先申請賣家審核。"),
]
yy = Inches(3.08)
for t, b in prep:
    tb = textbox(s, rx, yy, Inches(1.4), Inches(0.28))
    write(tb.text_frame, t, size=13, bold=True, color=ACCENT, first=True, space_after=0)
    tb = textbox(s, rx + Inches(1.5), yy, rw - Inches(1.5), Inches(0.6))
    write(tb.text_frame, b, size=12.5, color=MUTED, line=1.45, first=True, space_after=0)
    yy += Inches(0.76)

card(s, ML, Inches(5.36), CW, Inches(1.24), fill=SURFACE, line_color=LINE)
tb = textbox(s, ML + Inches(0.36), Inches(5.6), CW - Inches(0.72), Inches(0.8))
write(tb.text_frame,
      [("誠實揭露：", {"bold": True, "color": DANGER}),
       ("這一頁的每一條都是限制，我們沒有把它包裝成「新手保護」。"
        "它存在的理由是：沒有保證金的情況下，平台需要一個上界來承擔第一次違約的風險；"
        "違約累積 3 次就不能再開新池。你如果覺得這條不合理，開池前就該知道，"
        "而不是上線之後才發現。", {"color": MUTED})],
      size=13, line=1.55, first=True, space_after=0)

d.record("新賣家", "第一個池有上限。這是沒有保證金的代價，也是它的好處。",
         "我們不收保證金。取而代之的是把第一次違約的最大損失壓住 —— 限制只套用在你的第一個池。",
         ["第一個池上限：100 籤 / 100,000 點票收。完成第一個池就解除，之後沒有任何規模限制。",
          "開始要準備的東西：手上的實體卡（鑑定卡填機構與編號，生卡直接登記）／收件資訊由平台在買家申請出貨時提供／一個 LINE 或 Email 帳號，不用先申請賣家審核。",
          "誠實揭露：這一頁每一條都是限制，我們沒有包裝成「新手保護」。它存在的理由是沒有保證金時平台需要一個上界；違約累積 3 次就不能再開新池。"])

# ══════════════════════════════════════════════════════════════════════
# 13 現在的狀態
# ══════════════════════════════════════════════════════════════════════
s, y = d.new("現況", "平台還在測試階段，尚未正式營運。",
             "這一頁不是行銷。你如果現在要把庫存投進來，該先知道下面三件事。")

states = [
    ("測試階段", "尚未正式營運。站上目前陳列的池是示範資料，不是真實交易。"
                 "上線時程確定後會直接通知已登記的賣家。", WARN),
    ("PSA 編號自動查證", "還在等 PSA 核准 API。目前鑑定卡標成「暫未驗證」，"
                         "但不擋開池 —— 編號的唯一性約束照常生效。", WARN),
    ("已經可以用的", "抽成 0%、逐筆結算、買回價自訂、編號唯一性、"
                     "commit-reveal 公平性、經濟護欄 —— 本簡報前面幾頁的機制"
                     "都已經實作完成，畫面就是實際跑出來的。", OK),
]
cw3 = (CW - Inches(0.44)) / 3
for i, (t, b, c) in enumerate(states):
    x = ML + int(i * (cw3 + Inches(0.22)))
    card(s, x, Inches(2.62), int(cw3), Inches(2.72))
    rule(s, x + Inches(0.32), Inches(2.62), Inches(0.8), color=c, thickness=Pt(3))
    tb = textbox(s, x + Inches(0.32), Inches(3.02), int(cw3) - Inches(0.64), Inches(0.32))
    write(tb.text_frame, t, size=16, bold=True, color=TEXT, first=True, space_after=0)
    tb = textbox(s, x + Inches(0.32), Inches(3.5), int(cw3) - Inches(0.64), Inches(1.6))
    write(tb.text_frame, b, size=12.5, color=MUTED, line=1.5, first=True, space_after=0)

card(s, ML, Inches(5.6), CW, Inches(1.0), fill=SURFACE2, line_color=None)
tb = textbox(s, ML + Inches(0.36), Inches(5.84), CW - Inches(0.72), Inches(0.56))
write(tb.text_frame,
      [("本簡報沒有出現任何使用者數、交易量或成長率。", {"bold": True, "color": TEXT}),
       ("原因很簡單：平台還沒營運，那些數字現在都不存在。"
        "會出現的每一個數字，都是系統裡設定好的規則。", {"color": MUTED})],
      size=13.5, line=1.5, first=True, space_after=0)

d.record("現況", "平台還在測試階段，尚未正式營運。",
         "這一頁不是行銷。你如果現在要把庫存投進來，該先知道下面三件事。",
         ["測試階段 —— 尚未正式營運。站上目前陳列的池是示範資料，不是真實交易。上線時程確定後會直接通知已登記的賣家。",
          "PSA 編號自動查證 —— 還在等 PSA 核准 API。目前鑑定卡標成「暫未驗證」，但不擋開池，編號的唯一性約束照常生效。",
          "已經可以用的 —— 抽成 0%、逐筆結算、買回價自訂、編號唯一性、commit-reveal 公平性、經濟護欄都已實作完成，畫面就是實際跑出來的。",
          "本簡報沒有出現任何使用者數、交易量或成長率：平台還沒營運，那些數字現在都不存在。"])

# ══════════════════════════════════════════════════════════════════════
# 14 客人從哪裡來（行銷計畫）
#
# 刻意排在「現況」之後：讀者先知道平台還沒營運、簡報裡沒有任何實績數字，
# 再看這一頁，它就讀成「我們知道現況，這是打算」，而不是推銷。放在 13 之前
# 會變成先賣願景、再補上「其實還沒上線」，那個順序會把前面十三頁的可信度賠掉。
#
# 這一頁一樣不准出現預估數字。
#
# 廣告那一條改寫過（2026-09 重查一手政策原文）。前一版寫「付費廣告投不了」，
# 那是**過度推論**：Meta「Online Gambling and Games」對博弈的定義是
# "any product or service where anything of monetary value is included as part
# of a method of entry and prize"，舉的例子是現金與加密貨幣；台灣確實在該類
# 的 unsupported markets 清單上，但**實體商品的抽選包不在那個定義的明文範圍**。
# 實務上同類平台（Clove 等）在 Meta 都有經營。所以真正的變數是送審時被歸類
# 成什麼，而那取決於素材怎麼拍 —— 把它寫成「投不了」既不準確，也等於先幫
# 合作對象放棄一條真的存在的路。
#
# 另外補上一個具體且可執行的要求：台灣投放需先完成廣告主驗證
# （詐欺犯罪危害防制條例，2025-01-06 起）。
#
# 語氣也重寫過。前一版九條裡有三條是「投不了」「不會辦的」「做不到的」，
# 在一份要說服店家的簡報裡讀起來像免責聲明。界線該守的還是守（該說的話沒有
# 刪掉，只是改成正面表述 —— 例如「聯名池沒有」寫成「你的池只有你的招牌」），
# 但一頁行銷計畫的主詞應該是「我們會做什麼」，不是「我們不做什麼」。
# ══════════════════════════════════════════════════════════════════════
MKT_SUB = ("「抽成 0%」講的是你的成本，這一頁講的是流量。它是一份計畫 —— "
           "沒有預估人數、沒有轉換率，平台還沒營運，那些數字給不出來。")
MKT_SRC = ("政策查證（2026-09）：Meta 以「投注與獎項含金錢價值」界定博弈，實體商品的抽選包不在明文範圍 —— "
           "是送審歸類問題，不是禁止。台灣投放另需廣告主驗證（2025-01 起）。")
s, y = d.new("客人從哪裡來", "你把庫存投進來，客人從哪裡來？", MKT_SUB)

channels = [
    ("社群", "內容打底，廣告加速", GOLD, [
        "廣告投得出去，變數在素材 —— 擋的是「拿錢換錢」，我們的獎品是實體卡；素材講卡況，不講中獎金額。",
        "台灣投放先過廣告主驗證 —— 2025 年 1 月起的法定要求，這關由平台跑完，你只要出素材。",
        "我們發的內容是驗算 —— 完抽後公開種子，任何人都能自己算一次。這是別家抄不走的素材。",
    ]),
    ("活動", "平台出版位與檔期", LILAC, [
        "首發池排程 —— 新賣家的第一個池排進大廳的固定版位，不必一開始就跟老池搶自然排序。",
        "主題檔期 —— 同一個系列或主題的池集中在同一週開，多個賣家共用一次流量。參不參加你決定。",
        "開箱合作 —— 平台幫你牽線創作者。業配一律標示清楚（公平交易法），標了反而更有人信。",
    ]),
    ("跟店家合作", "你的店與你的池互相帶客", OK, [
        "賣家頁是你的線上門面 —— 開過的池、已出貨卡片、平均出貨天數、糾紛率，跟你的簡介同一頁。",
        "池有公開網址 —— 印成 QR 貼在櫃檯、發進你自己的 LINE 群與社團。店裡看到卡，回家抽你的池。",
        "你的池只掛你的招牌 —— 一個池屬於一個賣家，客人在你的池頁面不會看到別人的名字。",
    ]),
]
# 三橫列而不是三直欄：每一條的正文都是一整句，直欄只剩 18 字寬會逼出三行、
# 行距一撞就疊在一起（實測過）。橫列一行放得下 50 字，正文才讀得順。
RX = ML + Inches(3.10)              # 右半（條目）起點
RW = CW - Inches(3.44)
for i, (t, tag, c, rows) in enumerate(channels):
    R = Inches(2.70) + int(i * Inches(1.34))
    card(s, ML, R, CW, Inches(1.22))
    rule(s, ML + Inches(0.30), R, Inches(0.8), color=c, thickness=Pt(3))
    tb = textbox(s, ML + Inches(0.32), R + Inches(0.26), Inches(2.6), Inches(0.32))
    write(tb.text_frame, t, size=16, bold=True, color=TEXT, first=True, space_after=0)
    tb = textbox(s, ML + Inches(0.32), R + Inches(0.66), Inches(2.6), Inches(0.26))
    write(tb.text_frame, tag, size=11.5, bold=True, color=c, first=True, space_after=0)
    for j, body in enumerate(rows):
        lead, rest = body.split(" —— ", 1)
        tb = textbox(s, RX, R + Inches(0.22) + int(j * Inches(0.32)), RW, Inches(0.30))
        write(tb.text_frame,
              [(lead + " —— ", {"bold": True, "color": TEXT}), (rest, {"color": MUTED})],
              size=11.5, line=1.35, first=True, space_after=0)

caption(s, ML, Inches(6.66), CW, MKT_SRC)

# outline 的條目直接從 channels 長出來，不另外手打一份 —— 這一頁的資料只有一份。
d.record("客人從哪裡來", "你把庫存投進來，客人從哪裡來？", MKT_SUB,
         [f"{t}（{tag}）· {body}" if j == 0 else f"{t} · {body}"
          for t, tag, _c, rows in channels
          for j, body in enumerate(rows)] + [MKT_SRC])

# ══════════════════════════════════════════════════════════════════════
# 15 聯絡方式
# ══════════════════════════════════════════════════════════════════════
s, y = d.new("下一步", "想先開一個池試試看？",
             "第一個池上限 100 籤 / 100,000 點，沒有費用、不用保證金、不用審核排隊。")

card(s, ML, Inches(2.72), Inches(6.4), Inches(3.4), fill=SURFACE, line_color=LINE)
tb = textbox(s, ML + Inches(0.44), Inches(3.06), Inches(5.5), Inches(0.32))
write(tb.text_frame, "聯絡方式", size=16, bold=True, color=TEXT, first=True, space_after=0)
fields = ["聯絡人", "電話 / LINE", "Email", "網站"]
yy = Inches(3.62)
for f in fields:
    tb = textbox(s, ML + Inches(0.44), yy, Inches(1.4), Inches(0.28))
    write(tb.text_frame, f, size=12.5, color=MUTED, first=True, space_after=0)
    rule(s, ML + Inches(1.94), yy + Inches(0.26), Inches(4.0),
         color=RGBColor(0x3A, 0x36, 0x44))
    yy += Inches(0.62)

rx = ML + Inches(6.8)
rw = CW - Inches(6.8)
tb = textbox(s, rx, Inches(2.72), rw, Inches(0.32))
write(tb.text_frame, "回顧：為什麼是這裡", size=16, bold=True, color=TEXT, first=True,
      space_after=0)
recap = [
    ("抽成 0%", "票收扣掉你宣告的買回價，剩下全是你的"),
    ("逐筆結算", "一筆走完交付就入袋，不等整池"),
    ("買回價自己定", "絕對金額，開賣後鎖死，平台不調"),
    ("一卡不會被賣兩次", "資料庫層的唯一約束，護你的信譽"),
    ("公平性可驗算", "客訴有客觀答案，不用自證清白"),
]
yy = Inches(3.2)
for t, b in recap:
    pill(s, rx + Inches(0.02), yy + Inches(0.08), Inches(0.11), Inches(0.11), "", ACCENT, ACCENT, size=6)
    tb = textbox(s, rx + Inches(0.3), yy, Inches(2.1), Inches(0.28))
    write(tb.text_frame, t, size=13, bold=True, color=TEXT, first=True, space_after=0)
    tb = textbox(s, rx + Inches(0.3), yy + Inches(0.3), rw - Inches(0.3), Inches(0.28))
    write(tb.text_frame, b, size=11.5, color=MUTED, first=True, space_after=0)
    yy += Inches(0.68)

caption(s, ML, Inches(6.4), CW,
        "平台尚未正式營運。本簡報所有數字皆為系統設定值，非營運實績。")

d.record("下一步", "想先開一個池試試看？",
         "第一個池上限 100 籤 / 100,000 點，沒有費用、不用保證金、不用審核排隊。",
         ["聯絡方式（留空待填）：聯絡人 / 電話・LINE / Email / 網站。",
          "回顧：抽成 0% —— 票收扣掉你宣告的買回價，剩下全是你的。",
          "回顧：逐筆結算 —— 一筆走完交付就入袋，不等整池。",
          "回顧：買回價自己定 —— 絕對金額，開賣後鎖死，平台不調。",
          "回顧：一卡不會被賣兩次 —— 資料庫層的唯一約束，護你的信譽。",
          "回顧：公平性可驗算 —— 客訴有客觀答案，不用自證清白。",
          "平台尚未正式營運。本簡報所有數字皆為系統設定值，非營運實績。"])

# ══════════════════════════════════════════════════════════════════════
# ── 版面檢查：任何形狀跑出頁面就是錯 ────────────────────────────────
# 文字方塊的高度是「保留給文字的框」，實際文字可能更高，所以這一關只抓
# 幾何越界；文字有沒有溢出仍要靠 soffice 轉圖逐頁看。
def audit(prs):
    bad = []
    for i, sl in enumerate(prs.slides, 1):
        for sh in sl.shapes:
            L, T = sh.left, sh.top
            R, B = L + sh.width, T + sh.height
            if L < 0 or T < 0 or R > W + 1 or B > H + 1:
                bad.append((i, sh.shape_type, sh.name,
                            round(L / 914400, 2), round(T / 914400, 2),
                            round(R / 914400, 2), round(B / 914400, 2)))
    return bad


_bad = audit(d.prs)
for b in _bad:
    print(f"  越界　第 {b[0]} 頁 {b[2]}：{b[3]},{b[4]} → {b[5]},{b[6]}")
print(f"版面檢查：越界形狀 {len(_bad)} 個")

d.prs.save(OUT_PPTX)
print(f"寫出 {OUT_PPTX}（{d.n} 頁）")

# ── outline.md：跟簡報同一份資料 ──────────────────────────────────────
lines = [
    "# VaultDraw 賣家合作簡報 · 文字大綱",
    "",
    "> 這份大綱與 `VaultDraw-賣家合作簡報.pptx` 由 `build.py` 的同一份資料產生，",
    "> 兩邊內容必定一致。做網頁版請以本檔為準。",
    "",
    "**對象**：台灣的卡舖、卡商、個人大盤。",
    "**這份簡報要回答的唯一問題**：「我為什麼要在你這裡開池，而不是別家？」",
    "",
    "**硬性事實來源**：所有數字皆取自原始碼常數 —— "
    "`PLATFORM_FEE_RATE = 0`、`POOL_SHIP_DEADLINE_MS = 72h`、"
    "`POOL_INSPECT_MS = 7d`、`POOL_VAULT_ACCEPT_MS = 14d`、"
    "`BUYBACK_MIN = 10` / `BUYBACK_MAX = 10,000,000`、"
    "`POOL_DEFAULT_DAYS = 14` / `POOL_MAX_DAYS = 90`、"
    "`FLOOR_MINT = 100` / `FLOOR_THIN = 90` / `FLOOR_PREDATORY = 25`、"
    "`FIRST_POOL_TICKET_CAP = 100` / `FIRST_POOL_VALUE_CAP = 100,000`、"
    "`unique(grader, cert_no)`。**沒有任何使用者數、交易量或成長率** —— "
    "平台尚未營運，那些數字不存在。",
    "",
    "---",
    "",
]
for n, kicker, title, sub, bullets in d.outline:
    lines.append(f"## {n:02d} · {kicker}")
    lines.append("")
    lines.append(f"**{title}**")
    lines.append("")
    if sub:
        lines.append(sub)
        lines.append("")
    for b in bullets:
        lines.append(f"- {b}")
    lines.append("")

lines += [
    "---",
    "",
    "## 圖片來源",
    "",
    "全部由 headless Playwright 對本機 mock 模式（`VITE_API_URL= npm run dev`）",
    "以 1280 寬、`deviceScaleFactor: 2` 截取，放在 `assets/`：",
    "",
    "| 檔名 | 畫面 | 用在 |",
    "| --- | --- | --- |",
    "| `lobby.png` | 大廳首屏 | 01 封面背景 |",
    "| `x-tierbuy.png` | 開池表單・買回價（依賞別） | 06 |",
    "| `x-ret.png` | 池頁・保底回饋率 | 06 |",
    "| `x-certdup.png` | 開池表單・鑑定編號不能複製的錯誤列 | 07 |",
    "| `x-fair.png` | 池頁・Provably Fair 卡（commit hash 與 drand 輪次） | 08 |",
    "| `seller-new-filled-top.png` | 開池表單（套用範本後） | 09 |",
    "| `x-econ-over/warn/under/ok.png` | 即時試算的四種護欄狀態 | 10 |",
    "| `x-money.png` | 出貨與結算・可動用／保留中 | 11 |",
    "| `x-shiprow.png` | 出貨與結算・單筆出貨列與 72 小時倒數 | 11 |",
    "",
    "站上目前是示範資料，含有這類畫面的頁面都標了「示意畫面：站上目前為示範資料」。",
    "",
]
OUT_MD.write_text("\n".join(lines), encoding="utf-8")
print(f"寫出 {OUT_MD}")
